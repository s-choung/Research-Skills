#!/usr/bin/env python3
"""
PPTX Heavy Image Analyzer
- PPTX 파일의 모든 미디어를 슬라이드별로 매핑
- 2MB 이상 이미지를 추출 (EMF→PDF→PNG 변환 포함)
- 용량 내림차순 3열 그리드 HTML 생성
Usage: python analyze.py <pptx_path> [--threshold 2] [--output <dir>]
"""
import argparse
import zipfile
import os
import io
import tempfile
import shutil
import json
import subprocess
import base64
from collections import defaultdict
from pathlib import Path

try:
    from pptx import Presentation
    from PIL import Image
except ImportError:
    print("필요 패키지: pip install python-pptx Pillow")
    raise SystemExit(1)

SOFFICE_PATHS = [
    '/Applications/LibreOffice.app/Contents/MacOS/soffice',
    '/usr/bin/soffice',
    '/usr/local/bin/soffice',
]

def find_soffice():
    for p in SOFFICE_PATHS:
        if os.path.exists(p):
            return p
    return None

def convert_emf_to_png(emf_data, basename, tmpdir):
    soffice = find_soffice()
    if not soffice:
        return None, "N/A"

    emf_path = os.path.join(tmpdir, basename)
    pdf_path = emf_path.replace('.emf', '.pdf')
    png_path = emf_path.replace('.emf', '.png')

    with open(emf_path, 'wb') as f:
        f.write(emf_data)

    subprocess.run(
        [soffice, '--headless', '--convert-to', 'pdf', '--outdir', tmpdir, emf_path],
        capture_output=True, text=True, timeout=120,
        env={**os.environ, 'HOME': tmpdir}
    )

    if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 100:
        subprocess.run(
            ['magick', '-density', '150', pdf_path, '-strip', png_path],
            capture_output=True, text=True, timeout=120
        )
        if os.path.exists(png_path) and os.path.getsize(png_path) > 100:
            img = Image.open(png_path)
            return png_path, f"{img.size[0]}x{img.size[1]}"

    subprocess.run(
        [soffice, '--headless', '--convert-to', 'png', '--outdir', tmpdir, emf_path],
        capture_output=True, text=True, timeout=120,
        env={**os.environ, 'HOME': tmpdir}
    )
    if os.path.exists(png_path) and os.path.getsize(png_path) > 100:
        img = Image.open(png_path)
        return png_path, f"{img.size[0]}x{img.size[1]}"

    return None, "N/A"


def analyze_pptx(pptx_path, threshold_mb=2.0):
    prs = Presentation(pptx_path)
    z = zipfile.ZipFile(pptx_path, 'r')

    media_sizes = {}
    for info in z.infolist():
        if 'media/' in info.filename:
            media_sizes[info.filename] = info.file_size

    slide_media_map = defaultdict(list)
    for i, slide in enumerate(prs.slides, 1):
        for rel in slide.part.rels.values():
            target = rel.target_ref
            if 'media/' in str(target):
                parts = str(target).split('media/')
                mname = 'ppt/media/' + parts[-1] if len(parts) > 1 else str(target)
                if mname not in media_sizes:
                    for mf in media_sizes:
                        if mf.endswith(os.path.basename(str(target))):
                            mname = mf
                            break
                slide_media_map[mname].append(i)

    slide_labels = {}
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t and len(t) > 1:
                        texts.append(t)
        slide_labels[i] = texts[:3]

    threshold = int(threshold_mb * 1024 * 1024)
    big_media = sorted(
        [(n, s) for n, s in media_sizes.items() if s >= threshold],
        key=lambda x: x[1], reverse=True
    )

    z.close()

    return {
        'total_slides': len(prs.slides),
        'total_media': len(media_sizes),
        'total_media_mb': sum(media_sizes.values()) / 1024 / 1024,
        'file_size_mb': os.path.getsize(pptx_path) / 1024 / 1024,
        'big_media': big_media,
        'slide_media_map': slide_media_map,
        'slide_labels': slide_labels,
        'media_sizes': media_sizes,
    }


def extract_and_convert(pptx_path, big_media, out_dir, tmpdir):
    z = zipfile.ZipFile(pptx_path, 'r')
    results = []

    for name, size in big_media:
        basename = os.path.basename(name)
        ext = os.path.splitext(basename)[1].lower()
        data = z.read(name)
        out_name = basename
        dims = "N/A"

        if ext == '.emf':
            png_path, dims = convert_emf_to_png(data, basename, tmpdir)
            if png_path:
                out_name = basename.replace('.emf', '_emf.png')
                shutil.copy2(png_path, os.path.join(out_dir, out_name))
            else:
                out_name = None
        elif ext in ('.tif', '.tiff'):
            img = Image.open(io.BytesIO(data))
            dims = f"{img.size[0]}x{img.size[1]}"
            out_name = basename.rsplit('.', 1)[0] + '_tif.png'
            img.save(os.path.join(out_dir, out_name), format='PNG')
        elif ext == '.wdp':
            out_name = None
            dims = "HD Photo"
        else:
            out_path = os.path.join(out_dir, basename)
            with open(out_path, 'wb') as f:
                f.write(data)
            try:
                img = Image.open(io.BytesIO(data))
                dims = f"{img.size[0]}x{img.size[1]}"
            except:
                pass

        results.append({
            'zip_name': name,
            'basename': basename,
            'out_name': out_name,
            'size_mb': size / 1024 / 1024,
            'ext': ext,
            'dims': dims,
        })

    z.close()
    return results


def size_class(mb):
    if mb >= 10:
        return 'huge'
    if mb >= 5:
        return 'big'
    if mb >= 3:
        return 'med'
    return 'low'


def size_color_class(mb):
    if mb >= 5:
        return ''
    if mb >= 3:
        return ' med'
    return ' low'


def ext_class(ext):
    if ext == '.emf':
        return ' emf'
    if ext in ('.tif', '.tiff'):
        return ' tif'
    return ''


def generate_html(analysis, images, out_dir, pptx_name):
    slide_media_map = analysis['slide_media_map']
    slide_labels = analysis['slide_labels']

    total_big_mb = sum(img['size_mb'] for img in images)

    cards = []
    for rank, img in enumerate(images, 1):
        slides = sorted(set(slide_media_map.get(img['zip_name'], [])))
        slide_strs = []
        for sn in slides:
            labels = slide_labels.get(sn, [])
            lbl = labels[0][:40] if labels else '(no text)'
            slide_strs.append(f"S{sn}: {lbl}")
        slide_text = ' &nbsp;|&nbsp; '.join(slide_strs) if slide_strs else '(unmapped)'

        if img['out_name']:
            thumb = f'<img class="thumb" src="{img["out_name"]}" loading="lazy">'
        else:
            thumb = '<div class="thumb no-preview">미리보기 불가</div>'

        sc = size_class(img['size_mb'])
        card_cls = 'emf' if img['ext'] == '.emf' else ('huge' if sc == 'huge' else '')

        ext_badge = ''
        if img['ext'] in ('.emf', '.tif', '.tiff', '.wdp'):
            ext_badge = f'<span class="ext{ext_class(img["ext"])}">{img["ext"].upper().strip(".")}</span>'

        cards.append(f'''<div class="card-wrap"><span class="rank">#{rank}</span>
<div class="card {card_cls}">
{thumb}
<div class="info">
<div class="row"><span class="name">{img['basename']}{ext_badge}</span><span class="size{size_color_class(img['size_mb'])}">{img['size_mb']:.2f} MB</span></div>
<div class="dims">{img['dims']}</div>
<div class="slides">{slide_text}</div>
</div></div></div>''')

    html = f'''<!DOCTYPE html>
<html lang="ko">
<head>
<meta charset="UTF-8">
<title>{pptx_name} — Heavy Images</title>
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{background:#111;color:#ddd;font-family:-apple-system,BlinkMacSystemFont,'Pretendard',sans-serif;padding:24px}}
h1{{text-align:center;font-size:1.4rem;font-weight:600;margin-bottom:6px;color:#fff}}
.sub{{text-align:center;font-size:.85rem;color:#888;margin-bottom:24px}}
.grid{{display:grid;grid-template-columns:repeat(3,1fr);gap:16px;max-width:1400px;margin:0 auto}}
.card{{background:#1a1a1a;border-radius:10px;overflow:hidden;border:1px solid #2a2a2a;transition:border-color .2s}}
.card:hover{{border-color:#555}}
.card.emf{{border-color:#c2420088}}
.card.huge{{border-color:#ef444488}}
.thumb{{width:100%;aspect-ratio:4/3;object-fit:contain;background:#0a0a0a;display:block}}
.no-preview{{display:flex;align-items:center;justify-content:center;color:#555;font-size:.8rem}}
.info{{padding:10px 12px 12px}}
.row{{display:flex;justify-content:space-between;align-items:baseline;margin-bottom:4px;gap:8px}}
.name{{font-size:.8rem;font-weight:600;color:#eee;font-family:'SF Mono',Menlo,monospace;white-space:nowrap;overflow:hidden;text-overflow:ellipsis}}
.size{{font-size:.85rem;font-weight:700;color:#f87171;font-family:'SF Mono',Menlo,monospace;white-space:nowrap}}
.size.med{{color:#fb923c}}
.size.low{{color:#facc15}}
.dims{{font-size:.72rem;color:#888;font-family:'SF Mono',Menlo,monospace}}
.ext{{display:inline-block;font-size:.65rem;font-weight:600;padding:1px 6px;border-radius:4px;background:#333;color:#aaa;margin-left:4px}}
.ext.emf{{background:#7f1d1d;color:#fca5a5}}
.ext.tif{{background:#713f12;color:#fde68a}}
.slides{{font-size:.72rem;color:#9ca3af;margin-top:4px;line-height:1.4}}
.rank{{position:absolute;top:8px;left:8px;background:#000c;color:#fff;font-size:.7rem;font-weight:700;padding:2px 7px;border-radius:4px;font-family:'SF Mono',Menlo,monospace}}
.card-wrap{{position:relative}}
.legend{{max-width:1400px;margin:24px auto 0;display:flex;gap:16px;justify-content:center;font-size:.75rem;color:#888}}
.legend span{{display:flex;align-items:center;gap:4px}}
.dot{{width:10px;height:10px;border-radius:50%;display:inline-block}}
.dot.r{{background:#f87171}}.dot.o{{background:#fb923c}}.dot.y{{background:#facc15}}
</style>
</head>
<body>
<h1>{pptx_name} — Heavy Images ({len(images)})</h1>
<div class="sub">
  {analysis['total_slides']} slides &middot;
  {analysis['total_media']} media files &middot;
  total {analysis['file_size_mb']:.1f} MB &middot;
  shown: {total_big_mb:.1f} MB in {len(images)} images
</div>
<div class="legend">
  <span><span class="dot r"></span> 5 MB+</span>
  <span><span class="dot o"></span> 3-5 MB</span>
  <span><span class="dot y"></span> 2-3 MB</span>
</div>
<br>
<div class="grid">
{"".join(cards)}
</div>
</body>
</html>'''

    html_path = os.path.join(out_dir, 'heavy_images.html')
    with open(html_path, 'w', encoding='utf-8') as f:
        f.write(html)
    return html_path


def main():
    parser = argparse.ArgumentParser(description='PPTX Heavy Image Analyzer')
    parser.add_argument('pptx', help='Path to .pptx file')
    parser.add_argument('--threshold', type=float, default=2.0, help='Minimum MB to include (default: 2)')
    parser.add_argument('--output', help='Output directory (default: /tmp/pptx_heavy_<name>)')
    args = parser.parse_args()

    pptx_path = os.path.abspath(args.pptx)
    if not os.path.exists(pptx_path):
        print(f"파일 없음: {pptx_path}")
        raise SystemExit(1)

    pptx_name = Path(pptx_path).stem
    out_dir = args.output or f'/tmp/pptx_heavy_{pptx_name[:30]}'
    os.makedirs(out_dir, exist_ok=True)
    tmpdir = tempfile.mkdtemp(prefix='pptx_emf_')

    print(f"분석: {pptx_name}")
    print(f"출력: {out_dir}")

    analysis = analyze_pptx(pptx_path, args.threshold)
    print(f"  {analysis['total_slides']} slides, {analysis['total_media']} media, {analysis['file_size_mb']:.1f} MB")
    print(f"  {args.threshold}MB+ 이미지: {len(analysis['big_media'])}개")

    if not analysis['big_media']:
        print(f"  {args.threshold}MB 이상 이미지 없음!")
        shutil.rmtree(tmpdir, ignore_errors=True)
        return

    print("이미지 추출/변환 중...")
    images = extract_and_convert(pptx_path, analysis['big_media'], out_dir, tmpdir)

    print("HTML 생성 중...")
    html_path = generate_html(analysis, images, out_dir, pptx_name)

    shutil.rmtree(tmpdir, ignore_errors=True)

    print(f"\n완료!")
    print(f"  HTML: {html_path}")
    print(f"  이미지: {out_dir}/")
    print(f"  큰 이미지 {len(images)}개, 합계 {sum(i['size_mb'] for i in images):.1f} MB")

    return html_path


if __name__ == '__main__':
    main()
