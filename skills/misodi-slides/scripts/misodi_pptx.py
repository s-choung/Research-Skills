# -*- coding: utf-8 -*-
"""misodi_pptx — 미소디(MISODI) 스타일 PPTX 빌더 라이브러리.

style-spec.md(SSOT)의 실측값을 코드 상수로 고정한다.
사용법: 이 모듈을 import해서 덱 빌드 스크립트를 작성한다 (example_build.py 참고).

핵심 API
--------
deck = MisodiDeck("out.pptx")
s = deck.add_title_slide(kicker=..., title_runs=[...], date=..., leader=...)
s = deck.add_content_slide("제목", subtitle="부제", page=(2, 16))
deck.panel(s, x, y, w, h, color="navy", tab="탭 제목", fill_white=True)
deck.card(s, x, y, w, h, color="navy")
deck.chip(s, x, y, w, "칩 텍스트", color="orange")
deck.summary_band(s, y, runs=[...])           # 골드 요약 밴드
deck.navy_footer(s, runs=[...])               # 네이비 푸터 밴드
deck.text(s, x, y, w, h, runs=[...], size=16, align="l")
deck.add_fade_animation(s, shapes=[...])      # 블록 fade-in 500ms/stagger 300ms
deck.add_closing_slide()
deck.save()                                    # 저장 시 모든 슬라이드에 fade 전환 주입
"""
import os
import copy
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ---------- 상수 (style-spec.md 실측값) ----------
CM = 360000  # EMU per cm
SLIDE_W = 12192000
SLIDE_H = 6858000

COLORS = {
    "navy": "0F0F70",
    "gold": "C5A86F",
    "gold_dark": "9B7C3F",
    "gold_dark2": "B6924A",
    "gold_light": "D5C097",
    "parchment": "DCDAB2",
    "orange": "FA901E",
    "orange2": "FC9828",
    "gray": "888888",
    "silver": "B5B6B6",
    "periwinkle": "626DAE",
    "periwinkle2": "475DA3",
    "periwinkle_light": "8787B8",
    "shadow_navy": "324A88",
    "ink": "000000",
    "red": "FF0000",
    "title_blue": "8AB1E3",
    "title_peach": "FFC18C",
    "white": "FFFFFF",
}
# 섹션색 -> (테두리/칩 색, 내부 소제목 색, 내부 카드 테두리, 그림자 스트립 색+알파%)
SECTION = {
    "navy":   ("0F0F70", "0F0F70", "626DAE", ("324A88", 10)),
    "gold":   ("C5A86F", "9B7C3F", "D5C097", ("DCDAB2", 31)),
    "gray":   ("888888", "595959", "B5B6B6", ("888888", 10)),
    "orange": ("FA9028", "FA901E", "FDC88F", ("FEE4C7", 31)),
}
FONT_XB = "나눔스퀘어 ExtraBold"
FONT_B = "나눔스퀘어 Bold"

# 라운드 코너 어휘 (adj/100000 = 반경비)
ADJ_CONTAINER = 6386
ADJ_CARD = 12808
ADJ_BAND = 23741
ADJ_TAB = 39062
ADJ_PILL = 50000

ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "assets")
HEADER_IMG = os.path.join(ASSETS, "header_band.png")
WAVE_BG = os.path.join(ASSETS, "wave_bg.jpg")

# fade 스펙
FADE_DUR = 500       # ms, 요소 등장
FADE_STAGGER = 300   # ms
TRANSITION_DUR = 700 # ms, 슬라이드 전환


def cm(v):
    return Emu(int(round(v * CM)))


# ---------- low-level XML helpers ----------

def _set_adjs(shape, adjs):
    """프리셋 지오메트리의 adj 값들을 정확히 지정. adjs = {"adj": v} 또는 {"adj1": v, "adj2": v}."""
    sp = shape._element
    prst = sp.find(".//" + qn("a:prstGeom"))
    for child in list(prst):
        prst.remove(child)
    avLst = etree.SubElement(prst, qn("a:avLst"))
    for name, val in adjs.items():
        gd = etree.SubElement(avLst, qn("a:gd"))
        gd.set("name", name)
        gd.set("fmla", "val %d" % val)


def _set_round(shape, adj):
    """roundRect의 adj 값을 정확히 지정."""
    _set_adjs(shape, {"adj": adj})


def _solid_fill(shape, hexcolor, alpha_pct=None):
    """솔리드 채움 + 선택적 알파(불투명도 %)."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(hexcolor)
    if alpha_pct is not None:
        srgb = shape.fill.fore_color._xFill.find(".//" + qn("a:srgbClr"))
        a = etree.SubElement(srgb, qn("a:alpha"))
        a.set("val", str(int(alpha_pct * 1000)))


def _grad_fill_vertical(shape, stops):
    """세로 그라데이션. stops = [(pos_pct, hex, alpha_pct|None), ...] 위->아래."""
    spPr = shape._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill"):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    grad = etree.SubElement(spPr, qn("a:gradFill"))
    gsLst = etree.SubElement(grad, qn("a:gsLst"))
    for pos, hexc, alpha in stops:
        gs = etree.SubElement(gsLst, qn("a:gs"))
        gs.set("pos", str(int(pos * 1000)))
        c = etree.SubElement(gs, qn("a:srgbClr"))
        c.set("val", hexc)
        if alpha is not None:
            a = etree.SubElement(c, qn("a:alpha"))
            a.set("val", str(int(alpha * 1000)))
    lin = etree.SubElement(grad, qn("a:lin"))
    lin.set("ang", "5400000")  # 90deg = top->bottom
    lin.set("scaled", "1")
    # spPr 내 위치 보정: ln 앞에 와야 함
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        spPr.remove(grad)
        ln.addprevious(grad)


def _line(shape, hexcolor, width_pt, alpha_pct=None):
    shape.line.color.rgb = RGBColor.from_string(hexcolor)
    shape.line.width = Pt(width_pt)
    if alpha_pct is not None:
        srgb = shape.line.color._xFill.find(".//" + qn("a:srgbClr"))
        a = etree.SubElement(srgb, qn("a:alpha"))
        a.set("val", str(int(alpha_pct * 1000)))


def _no_line(shape):
    shape.line.fill.background()


def _runs(paragraph, runs, default_size=16, align="l"):
    """runs = [(text, dict)] / dict keys: size, color, xb(ExtraBold), b(bold flag)"""
    paragraph.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}[align]
    for text, opt in runs:
        r = paragraph.add_run()
        r.text = text
        f = r.font
        f.size = Pt(opt.get("size", default_size))
        f.bold = opt.get("b", True)
        name = FONT_XB if opt.get("xb") else FONT_B
        f.name = name
        # ea(한글) 폰트도 동일 지정
        rPr = r._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", name)
        f.color.rgb = RGBColor.from_string(opt.get("color", "000000"))


# ---------- Deck ----------

class MisodiDeck:
    def __init__(self, out_path):
        self.out_path = out_path
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self._blank = self.prs.slide_layouts[6]
        self._anim_shapes = {}  # slide -> [shape,...]

    # ----- 슬라이드 -----
    def _new_slide(self):
        return self.prs.slides.add_slide(self._blank)

    def add_title_slide(self, kicker, title_runs, date, leader, leader_name=""):
        """타이틀: 풀블리드 wave 배경 + 좌측 정렬 텍스트 컬럼."""
        s = self._new_slide()
        s.shapes.add_picture(WAVE_BG, 0, 0, SLIDE_W, SLIDE_H)
        self.text(s, 2.16, 2.19, 14.0, 0.78, [(kicker, {"xb": True, "size": 16, "color": "FFFFFF"})])
        tb = self.text(s, 2.16, 4.64, 18.0, 6.77, title_runs, size=31)
        tb.text_frame.word_wrap = True
        self.text(s, 2.16, 12.98, 6.0, 0.81, [(date, {"xb": True, "size": 18, "color": "FFFFFF"})])
        t1 = self.text(s, 2.16, 15.92, 4.2, 0.81, [(leader, {"xb": True, "size": 20, "color": "FFFFFF"})])
        t1.text_frame.word_wrap = False
        ln = s.shapes.add_connector(1, cm(6.3), cm(15.81), cm(6.3), cm(16.66))
        _line(ln, "E0E0E0", 1.5)
        if leader_name:
            t2 = self.text(s, 6.7, 15.92, 5.0, 0.81, [(leader_name, {"xb": True, "size": 18, "color": "FFFFFF"})])
            t2.text_frame.word_wrap = False
        return s

    def add_content_slide(self, title, subtitle="", page=None, page_note=""):
        """헤더 밴드(이미지+오버레이+오렌지띠) + 제목 + 페이지 번호."""
        s = self._new_slide()
        # 원본 z-order 그대로: 오렌지 띠(맨 아래) -> 웨이브 이미지(띠 상부를 덮음) -> 다크 오버레이
        # 1) 오렌지 띠: y=1.74 h=0.36이지만 이미지(h=1.93)에 덮여 0.17cm만 보임
        strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, cm(1.74), SLIDE_W, cm(0.36))
        _solid_fill(strip, COLORS["orange2"], 88)
        _no_line(strip)
        # 2) 웨이브 이미지 스트립
        s.shapes.add_picture(HEADER_IMG, 0, 0, SLIDE_W, cm(1.93))
        # 3) 다크 그라데이션 오버레이: 좌측 불투명 다크그레이(#3A3939 = E7E6E6 lumMod 25%)
        #    -> 우측 black@59%. 원본 실측 (-0.01,-0.70) 33.88x1.96cm, ang=0
        ov = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cm(-0.01), cm(-0.70), cm(33.88), cm(1.96))
        spPr = ov._element.spPr
        ov.fill.solid()
        for el in spPr.findall(qn("a:solidFill")):
            spPr.remove(el)
        grad = etree.SubElement(spPr, qn("a:gradFill"))
        gsLst = etree.SubElement(grad, qn("a:gsLst"))
        for pos, hexc, alpha in [(0, "3A3939", None), (100, "000000", 59)]:
            gs = etree.SubElement(gsLst, qn("a:gs"))
            gs.set("pos", str(pos * 1000))
            c = etree.SubElement(gs, qn("a:srgbClr"))
            c.set("val", hexc)
            if alpha is not None:
                a = etree.SubElement(c, qn("a:alpha"))
                a.set("val", str(alpha * 1000))
        lin = etree.SubElement(grad, qn("a:lin"))
        lin.set("ang", "0")
        lin.set("scaled", "0")
        ln_el = spPr.find(qn("a:ln"))
        if ln_el is not None:
            spPr.remove(grad)
            ln_el.addprevious(grad)
        _no_line(ov)
        # 4) 제목
        runs = [(title, {"xb": True, "size": 28, "color": "FFFFFF"})]
        if subtitle:
            runs.append(("  -  " + subtitle, {"xb": True, "size": 22, "color": "FFFFFF"}))
        self.text(s, 0.80, 0.50, 28.5, 1.02, runs)
        # 5) 페이지 번호
        if page:
            ptxt = "%d/%d" % page + ((" (" + page_note + ")") if page_note else "")
            self.text(s, 29.5, 0.55, 4.0, 0.8, [(ptxt, {"size": 12, "color": "FFFFFF", "b": False})], align="r")
        return s

    def add_closing_slide(self, message="경청해 주셔서 감사합니다.", sub_runs=None):
        s = self._new_slide()
        s.shapes.add_picture(WAVE_BG, 0, 0, SLIDE_W, SLIDE_H)
        t = self.text(s, 4.0, 8.2, 25.87, 1.5, [(message, {"xb": True, "size": 28, "color": "FFFFFF"})], align="c")
        t.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        if sub_runs:
            self.text(s, 4.0, 16.5, 25.87, 1.2, sub_runs, size=12, align="c")
        return s

    # ----- 컴포넌트 -----
    def shadow_strip(self, s, x, y, w, color="navy", h=3.0, corner_cm=None):
        """패널 뒤 깔개. 위쪽 모서리만 컨테이너 곡률에 맞춰 둥글게 (뾰족 튀어나옴 방지)."""
        hexc, alpha = SECTION[color][3]
        sh = s.shapes.add_shape(MSO_SHAPE.ROUND_2_SAME_RECTANGLE,
                                cm(x), cm(y), cm(w), cm(h))
        # round2SameRect: corner radius = adj1/100000 * min(w,h)
        r = corner_cm if corner_cm else ADJ_CONTAINER / 100000.0 * min(w, h)
        adj1 = int(round(r / min(w, h) * 100000))
        _set_adjs(sh, {"adj1": max(0, min(50000, adj1)), "adj2": 0})
        _solid_fill(sh, hexc, alpha)
        _no_line(sh)
        return sh

    def panel(self, s, x, y, w, h, color="navy", tab=None, tab_w=None, fill_white=True, shadow=True):
        """메인 패널: 그림자 스트립 + 컨테이너(adj=6386, 2.25pt) + 탭 칩(adj=39062, 겹침)."""
        line_c = SECTION[color][0]
        if shadow:
            # 스트립 위 모서리 반경 = 컨테이너 반경 (adj=6386 * 패널 짧은변)
            self.shadow_strip(s, x, y, w, color,
                              corner_cm=ADJ_CONTAINER / 100000.0 * min(w, h))
        cont = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cm(x), cm(y), cm(w), cm(h))
        _set_round(cont, ADJ_CONTAINER)
        if fill_white:
            _solid_fill(cont, "FFFFFF")
        else:
            cont.fill.background()
        _line(cont, line_c, 2.25)
        chip = None
        if tab:
            tw = tab_w if tab_w else min(11.0, w * 0.72)
            ch = 1.35
            cx = x + (w - tw) / 2
            cy = y - ch / 2  # 컨테이너 top에 칩 세로중심이 겹침 (시그니처)
            chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cm(cx), cm(cy), cm(tw), cm(ch))
            _set_round(chip, ADJ_TAB)
            _solid_fill(chip, line_c)
            _no_line(chip)
            tf = chip.text_frame
            tf.margin_top = tf.margin_bottom = 0
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            _runs(tf.paragraphs[0], [(tab, {"xb": True, "size": 20, "color": "FFFFFF"})], align="c")
        return cont, chip

    def card(self, s, x, y, w, h, color="navy", width_pt=1.0, adj=ADJ_CARD):
        """내부 카드: noFill + 얇은 테두리."""
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cm(x), cm(y), cm(w), cm(h))
        _set_round(c, adj)
        c.fill.background()
        _line(c, SECTION[color][2], width_pt)
        return c

    def grad_card(self, s, x, y, w, h, color="orange"):
        """그라데이션 카드: 위 틴트 -> 아래 흰색, 테두리 2.25pt 섹션색."""
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cm(x), cm(y), cm(w), cm(h))
        _set_round(c, ADJ_CONTAINER)
        if color == "orange":
            _grad_fill_vertical(c, [(0, "FDC88F", None), (34, "FEE4C7", None), (100, "FFFFFF", None)])
        else:
            _grad_fill_vertical(c, [(0, "8787B8", None), (34, "C3C3DC", None), (100, "FFFFFF", None)])
        _line(c, SECTION[color][0], 2.25)
        return c

    def chip(self, s, x, y, w, label, color="orange", h=1.0, size=16):
        """작은 칩 (adj=39062)."""
        line_c = SECTION[color][0] if color in SECTION else COLORS[color]
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cm(x), cm(y), cm(w), cm(h))
        _set_round(c, ADJ_TAB)
        _solid_fill(c, line_c)
        _no_line(c)
        tf = c.text_frame
        tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _runs(tf.paragraphs[0], [(label, {"xb": True, "size": size, "color": "FFFFFF"})], align="c")
        return c

    def pill(self, s, x, y, w, label, color_hex="475DA3", h=1.15, size=18):
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cm(x), cm(y), cm(w), cm(h))
        _set_round(c, ADJ_PILL)
        _solid_fill(c, color_hex)
        _no_line(c)
        tf = c.text_frame
        tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _runs(tf.paragraphs[0], [(label, {"xb": True, "size": size, "color": "FFFFFF"})], align="c")
        return c

    def summary_band(self, s, y, runs, x=9.29, w=15.28, h=2.46):
        """골드 요약 밴드 (adj=23741, 2.25pt #C5A86F)."""
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cm(x), cm(y), cm(w), cm(h))
        _set_round(c, ADJ_BAND)
        _solid_fill(c, "FFFFFF")
        _line(c, COLORS["gold"], 2.25)
        tf = c.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _runs(tf.paragraphs[0], runs, default_size=18, align="c")
        return c

    def navy_footer(self, s, runs, h=2.47):
        """슬라이드 최하단 네이비 풀폭 밴드."""
        y = SLIDE_H - cm(h)
        c = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, SLIDE_W, cm(h))
        _solid_fill(c, COLORS["navy"])
        _no_line(c)
        tf = c.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, (txt, opt) in enumerate(runs):
            opt.setdefault("color", "FFFFFF")
            opt.setdefault("size", 20)
        _runs(tf.paragraphs[0], runs, default_size=20, align="c")
        return c

    def divider_v(self, s, x, y, h):
        """중앙 세로 구분선: 그라데이션 [흰->네이비->흰] 2.25pt."""
        ln = s.shapes.add_connector(1, cm(x), cm(y), cm(x), cm(y + h))
        ln.line.width = Pt(2.25)
        lnEl = ln._element.spPr.find(qn("a:ln"))
        grad = etree.SubElement(lnEl, qn("a:gradFill"))
        gsLst = etree.SubElement(grad, qn("a:gsLst"))
        for pos, hexc in [(0, "FFFFFF"), (50, "0F0F70"), (100, "FFFFFF")]:
            gs = etree.SubElement(gsLst, qn("a:gs"))
            gs.set("pos", str(pos * 1000))
            c = etree.SubElement(gs, qn("a:srgbClr"))
            c.set("val", hexc)
        lin = etree.SubElement(grad, qn("a:lin"))
        lin.set("ang", "5400000")
        return ln

    def text(self, s, x, y, w, h, runs, size=16, align="l", anchor="t"):
        """noFill 텍스트박스. runs=[(text, opt)]"""
        tb = s.shapes.add_textbox(cm(x), cm(y), cm(w), cm(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(36000)
        tf.margin_top = tf.margin_bottom = Emu(18000)
        if anchor == "m":
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _runs(tf.paragraphs[0], runs, default_size=size, align=align)
        return tb

    def bullet_text(self, s, x, y, w, h, items, size=16):
        """여러 줄 본문. items = [runs, runs, ...]"""
        tb = s.shapes.add_textbox(cm(x), cm(y), cm(w), cm(h))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, runs in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            _runs(p, runs, default_size=size, align="l")
        return tb

    # ----- 애니메이션 -----
    def add_fade_animation(self, s, shapes):
        """지정한 도형들에 fade-in(500ms) entrance를 읽기 순서대로 300ms stagger로 주입.
        슬라이드 표시 직후 자동 시작. HTML 버전과 동일 타이밍."""
        self._anim_shapes[s.slide_id] = list(shapes)

    def _write_timing(self, s, shapes):
        spids = [sh.shape_id for sh in shapes]
        nsmap_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        P = "{%s}" % nsmap_p
        sld = s._element
        old = sld.find(P + "timing")
        if old is not None:
            sld.remove(old)

        def E(parent, tag, **attrs):
            el = etree.SubElement(parent, P + tag)
            for k, v in attrs.items():
                el.set(k, str(v))
            return el

        timing = etree.SubElement(sld, P + "timing")
        tnLst = E(timing, "tnLst")
        par_root = E(tnLst, "par")
        ctn_root = E(par_root, "cTn", id=1, dur="indefinite", restart="never", nodeType="tmRoot")
        child_root = E(ctn_root, "childTnLst")
        seq = E(child_root, "seq", concurrent=1, nextAc="seek")
        ctn_main = E(seq, "cTn", id=2, dur="indefinite", nodeType="mainSeq")
        main_child = E(ctn_main, "childTnLst")

        # 단일 그룹: 첫 효과는 afterEffect(슬라이드 표시 후 자동), 나머지는 withEffect + delay
        gpar = E(main_child, "par")
        gctn = E(gpar, "cTn", id=3, fill="hold")
        gst = E(gctn, "stCondLst")
        E(gst, "cond", delay="indefinite")
        onb = E(gst, "cond", evt="onBegin", delay=0)
        E(onb, "tn", val=2)
        gchild = E(gctn, "childTnLst")
        ipar = E(gchild, "par")
        ictn = E(ipar, "cTn", id=4, fill="hold")
        ist = E(ictn, "stCondLst")
        E(ist, "cond", delay=0)
        ichild = E(ictn, "childTnLst")

        nid = 5
        for i, spid in enumerate(spids):
            epar = E(ichild, "par")
            ectn = E(epar, "cTn", id=nid, presetID=10, presetClass="entr",
                     presetSubtype=0, fill="hold", grpId=0,
                     nodeType=("afterEffect" if i == 0 else "withEffect"))
            nid += 1
            est = E(ectn, "stCondLst")
            E(est, "cond", delay=i * FADE_STAGGER)
            echild = E(ectn, "childTnLst")
            # set visibility
            setEl = E(echild, "set")
            cBhvr = E(setEl, "cBhvr")
            sctn = E(cBhvr, "cTn", id=nid, dur=1, fill="hold")
            nid += 1
            sst = E(sctn, "stCondLst")
            E(sst, "cond", delay=0)
            tgt = E(cBhvr, "tgtEl")
            E(tgt, "spTgt", spid=spid)
            attrLst = E(cBhvr, "attrNameLst")
            attr = etree.SubElement(attrLst, P + "attrName")
            attr.text = "style.visibility"
            toEl = E(setEl, "to")
            strVal = E(toEl, "strVal", val="visible")
            # animEffect fade
            ae = E(echild, "animEffect", transition="in", filter="fade")
            cBhvr2 = E(ae, "cBhvr")
            actn = E(cBhvr2, "cTn", id=nid, dur=FADE_DUR)
            nid += 1
            tgt2 = E(cBhvr2, "tgtEl")
            E(tgt2, "spTgt", spid=spid)

        # prevCondLst / nextCondLst (클릭 내비게이션)
        prevC = E(seq, "prevCondLst")
        pc = E(prevC, "cond", evt="onPrev", delay=0)
        E(pc, "tgtEl").append(etree.SubElement(pc.find(P + "tgtEl"), P + "sldTgt"))
        nextC = E(seq, "nextCondLst")
        nc = E(nextC, "cond", evt="onNext", delay=0)
        E(nc, "tgtEl").append(etree.SubElement(nc.find(P + "tgtEl"), P + "sldTgt"))

    def _write_transition(self, s):
        """fade 슬라이드 전환 (700ms)."""
        P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
        P14 = "{http://schemas.microsoft.com/office/powerpoint/2010/main}"
        MC = "{http://schemas.openxmlformats.org/markup-compatibility/2006}"
        sld = s._element
        for old in sld.findall(P + "transition"):
            sld.remove(old)
        for old in sld.findall(MC + "AlternateContent"):
            if old.find(".//" + P + "transition") is not None:
                sld.remove(old)
        # mc:AlternateContent로 p14:dur(ms 단위) + fallback spd 지정
        ac = etree.SubElement(sld, MC + "AlternateContent")
        choice = etree.SubElement(
            ac, MC + "Choice",
            nsmap={"p14": "http://schemas.microsoft.com/office/powerpoint/2010/main"})
        choice.set("Requires", "p14")
        tr1 = etree.SubElement(choice, P + "transition")
        tr1.set("spd", "slow")
        tr1.set(P14 + "dur", str(TRANSITION_DUR))
        etree.SubElement(tr1, P + "fade")
        fallback = etree.SubElement(ac, MC + "Fallback")
        tr2 = etree.SubElement(fallback, P + "transition")
        tr2.set("spd", "slow")
        etree.SubElement(tr2, P + "fade")
        # cSld, clrMapOvr 다음 위치로 이동
        sld.remove(ac)
        clr = sld.find(P + "clrMapOvr")
        anchor = clr if clr is not None else sld.find(P + "cSld")
        anchor.addnext(ac)

    def save(self):
        for s in self.prs.slides:
            if self._anim_shapes.get(s.slide_id):
                self._write_timing(s, self._anim_shapes[s.slide_id])
            self._write_transition(s)
        self.prs.save(self.out_path)
        return self.out_path
